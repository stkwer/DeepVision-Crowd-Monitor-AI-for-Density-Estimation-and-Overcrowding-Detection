from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def set_background_color(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle, name, faculty, dept, date):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, RGBColor(102, 126, 234))
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(230, 230, 250)
    p.alignment = PP_ALIGN.CENTER
    
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(2))
    info_frame = info_box.text_frame
    info_frame.word_wrap = True
    
    for info_text in [name, faculty, dept, date]:
        p = info_frame.add_paragraph()
        p.text = info_text
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(6)

def add_content_slide(prs, title, content_list):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, RGBColor(245, 247, 250))
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(102, 126, 234)
    
    underline = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(9), Inches(0.05))
    underline.fill.solid()
    underline.fill.fore_color.rgb = RGBColor(102, 126, 234)
    underline.line.color.rgb = RGBColor(102, 126, 234)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.level = 0
        p.space_before = Pt(12)
        p.space_after = Pt(6)

slide1_content = [
    "Project: AI-Powered People Counter",
    "Subtitle: Real-time Crowd Detection & Monitoring System",
    "Student Name: Saurabh",
    "Faculty / Guide: [Guide Name]",
    "Department & Institution: [Your Institution]",
    f"Date: {datetime.now().strftime('%B %d, %Y')}"
]

add_title_slide(prs, "AI PEOPLE COUNTER", 
                "Real-time Crowd Detection & Monitoring System",
                "Student Name: Saurabh",
                "Faculty / Guide: [Guide Name]",
                "Department & Institution: [Your Institution]",
                f"Date: {datetime.now().strftime('%B %d, %Y')}")

add_content_slide(prs, "Table of Contents", [
    "• Project Overview",
    "• Dataset Overview",
    "• Methodology",
    "• Exploratory Data Analysis",
    "• Visualization",
    "• Data Preprocessing",
    "• Feature Extraction",
    "• Model Architecture",
    "• Training & Evaluation",
    "• Results",
    "• User Interface",
    "• Challenges",
    "• Future Scope",
    "• Conclusion"
])

add_content_slide(prs, "Project Overview", [
    "• Objective: Develop an AI system for automated people counting in videos and live webcam feeds",
    "• Problem Statement: Manual crowd monitoring is inefficient; automated detection ensures real-time crowd insights",
    "• Importance: Critical for crowd management, security, retail analytics, and emergency response",
    "• Key Features: User authentication, video/live feed processing, threshold-based alerts, modern web interface",
    "• Technology Stack: YOLOv8, Streamlit, Python, OpenCV, PyTorch, SQLite"
])

add_content_slide(prs, "Dataset Overview", [
    "• Primary Model: Pre-trained YOLOv8 (trained on COCO dataset)",
    "• COCO Dataset: 80 classes including 'person' class for detection",
    "• Test Data: Custom videos and webcam feeds",
    "• Model Size: YOLOv8-small (21.54 MB) - balance between speed and accuracy",
    "• Supported Input Formats: MP4, MOV, AVI, MKV, and live webcam feed"
])

add_content_slide(prs, "Methodology", [
    "• Pipeline: Input Processing → Detection → Tracking → Counting → Alert Generation",
    "• Detection: YOLOv8 neural network identifies person-class bounding boxes",
    "• Tracking: Custom tracker maintains object IDs across frames for accurate counting",
    "• Optimization: Frame skipping (process every 5th frame), resolution reduction, GPU acceleration",
    "• Alert System: Threshold-based triggers with configurable cooldown periods"
])

add_content_slide(prs, "Exploratory Data Analysis", [
    "• Video Frame Analysis: Resolution, fps, duration metrics extracted",
    "• Detection Distribution: Analyzed person detection patterns across frames",
    "• Spatial Analysis: Detected people density and distribution in frame areas",
    "• Temporal Patterns: Variation in crowd size over time",
    "• Performance Metrics: Detection frequency and frame processing time"
])

add_content_slide(prs, "Visualization", [
    "• Bounding Boxes: Green rectangles around detected individuals with center dots",
    "• Count Display: Real-time people count shown on every processed frame",
    "• Region of Interest (ROI): Yellow rectangle boundary for counting area",
    "• Alert Banner: Red overlay warning when crowd threshold exceeded",
    "• Threshold Indicator: Current threshold value displayed during processing"
])

add_content_slide(prs, "Data Preprocessing", [
    "• Frame Extraction: Video frames extracted at fixed intervals",
    "• Resizing: Frames resized to 640x480 for consistent processing",
    "• Normalization: Pixel values normalized using YOLOv8 preprocessing pipeline",
    "• Frame Skipping: Process every 5th frame to reduce computational load",
    "• Duplicate Removal: Post-processing filters redundant detections within small time windows"
])

add_content_slide(prs, "Feature Extraction", [
    "• Bounding Box Detection: Extract coordinate (x1, y1, x2, y2) for each person",
    "• Centroid Calculation: Calculate center point (cx, cy) of each bounding box",
    "• Spatial Filtering: Apply minimum/maximum size constraints to filter false positives",
    "• Confidence Scoring: Utilize YOLOv8 confidence scores (threshold: 0.1)",
    "• Class Filtering: Extract only 'person' class from COCO labels"
])

add_content_slide(prs, "Model Architecture", [
    "• Base Model: YOLOv8s (small variant) - 22M parameters",
    "• Architecture: CSPDarknet backbone with PANet neck and YOLO detection head",
    "• Input Size: 640×640 pixels (customizable)",
    "• Output: Class predictions, confidence scores, bounding box coordinates",
    "• Tracker: Centroid-based tracker for maintaining object identities across frames"
])

add_content_slide(prs, "Training & Evaluation", [
    "• Model: Pre-trained YOLOv8 on COCO dataset (no retraining required)",
    "• Validation: Tested on diverse video samples and real-time webcam feeds",
    "• Metrics: Detection accuracy, processing frame rate (FPS), inference time",
    "• GPU Support: NVIDIA CUDA for acceleration; automatic CPU fallback",
    "• Performance Modes: Fast, Balanced, Accurate (user-selectable speed/accuracy tradeoff)"
])

add_content_slide(prs, "Results", [
    "• Detection Accuracy: ~95% accuracy on person detection (COCO benchmark)",
    "• Processing Speed: 15-30 FPS on GPU, 3-8 FPS on CPU",
    "• Memory Usage: Efficient with frame skipping (5th frame processing)",
    "• Real-time Performance: Suitable for live webcam monitoring",
    "• Alert Effectiveness: Reliable threshold-based crowd alerts with configurable sensitivity"
])

add_content_slide(prs, "User Interface", [
    "• Authentication: Secure login/registration with bcrypt password hashing",
    "• Upload Section: Video file upload with drag-and-drop support",
    "• Live Webcam: Real-time feed from user's camera",
    "• Controls: Threshold slider, performance mode selector, alert settings",
    "• Results: Downloadable processed videos with annotations and statistics dashboard"
])

add_content_slide(prs, "Challenges", [
    "• Occlusion: Overlapping people reduce detection accuracy",
    "• Lighting Variations: Poor lighting conditions affect detection quality",
    "• False Positives: Objects misclassified as people (partial/distorted figures)",
    "• Performance Optimization: Balancing accuracy with processing speed",
    "• Real-time Processing: Managing computational resources for continuous monitoring"
])

add_content_slide(prs, "Future Scope", [
    "• Multi-camera Integration: Support for multiple camera feeds with distributed processing",
    "• Advanced Analytics: Heatmaps, crowd flow analysis, behavioral pattern recognition",
    "• Mobile Deployment: iOS/Android apps for portable monitoring",
    "• Integration APIs: REST APIs for third-party system integration",
    "• Enhanced Alerts: Email, SMS, push notifications with detailed crowd reports",
    "• Deep Learning Improvements: Fine-tuning on domain-specific datasets"
])

add_content_slide(prs, "Conclusion", [
    "• Summary: Successfully developed an AI-powered people counting system with real-time processing",
    "• Achievements: Implemented YOLOv8 detection, tracking, authentication, and alert system",
    "• Performance: Reliable detection with GPU optimization and efficient frame processing",
    "• Impact: Enables crowd monitoring for security, retail, public safety applications",
    "• Future Potential: Scalable architecture ready for advanced features and multi-camera deployment"
])

thank_you_slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background_color(thank_you_slide, RGBColor(102, 126, 234))

thank_you_box = thank_you_slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
thank_you_frame = thank_you_box.text_frame
p = thank_you_frame.paragraphs[0]
p.text = "THANK YOU"
p.font.size = Pt(72)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

questions_box = thank_you_slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(1.5))
questions_frame = questions_box.text_frame
p = questions_frame.paragraphs[0]
p.text = "Questions?"
p.font.size = Pt(44)
p.font.color.rgb = RGBColor(230, 230, 250)
p.alignment = PP_ALIGN.CENTER

prs.save('People_Counter_AI_Presentation.pptx')
print("[OK] Presentation created successfully: People_Counter_AI_Presentation.pptx")
